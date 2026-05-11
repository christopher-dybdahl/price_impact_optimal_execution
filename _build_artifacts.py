"""Build the LaTeX-include tables and the slides.pptx from saved/ artifacts.

Run AFTER the notebook has finished executing and saved/<name>/ directories
are populated.  Outputs:

    final_project/tables/headline_metrics.tex
    final_project/tables/tca_summary.tex
    final_project/tables/sensitivity.tex
    final_project/slides.pptx

Also runs the sensitivity sweep (wrong-model, signal delay, rho, H) that
the notebook does not run, using the same modular pipeline.
"""
from __future__ import annotations

import sys
from pathlib import Path

import numpy as np
import pandas as pd

ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parent
SAVED = ROOT / "saved"
TABLES = ROOT / "tables"
TABLES.mkdir(exist_ok=True)

sys.path.insert(0, str(ROOT / "src"))
import price_impact as pi  # noqa: E402


# ---------------------------------------------------------------------------
# Load the per-run artifacts produced by the notebook's modular sweep.
# ---------------------------------------------------------------------------
def load_run(name: str) -> dict:
    rd = SAVED / name
    metrics = pd.read_csv(rd / "metrics.csv", index_col=0)["value"]
    tca = pd.read_csv(rd / "tca_summary.csv", index_col=0)["total"]
    daily = pd.read_csv(rd / "daily_pnl.csv", index_col=0)
    return {"name": name, "metrics": metrics, "tca": tca, "daily": daily}


HEADLINE_RUNS = ["ow_daily", "ow_multi", "afs_daily", "afs_multi"]


def headline_table(runs: dict[str, dict]) -> str:
    """Build the LaTeX body of the headline metrics table."""
    rows = []
    for n in HEADLINE_RUNS:
        m = runs[n]["metrics"]
        rows.append(
            {
                "config": n,
                "n_days": int(float(m["n_days"])),
                "sharpe_sim": float(m["sharpe_sim"]),
                "total_pnl_sim": float(m["total_pnl_sim"]),
                "impact_cost": float(m["total_impact_cost"]),
                "max_dd": float(m["max_drawdown"]),
                "max_impact": float(m["max_impact_dislocation"]),
                "win_rate": float(m["win_rate"]),
            }
        )
    df = pd.DataFrame(rows)
    body_lines = [
        r"\begin{tabular}{lrrrrrrr}",
        r"\toprule",
        (r"Config & Days & Sharpe & Net P\&L (\$) & "
         r"Impact cost (\$) & Max DD (\$) & Max $|\lambda\bar I|$ & Win \% \\"),
        r"\midrule",
    ]
    for r in df.itertuples(index=False):
        cfg = r.config.replace("_", r"\_")
        body_lines.append(
            f"\\texttt{{{cfg}}} & {r.n_days} & {r.sharpe_sim:+.3f} & "
            f"{r.total_pnl_sim:,.0f} & {r.impact_cost:,.0f} & "
            f"{r.max_dd:,.0f} & {r.max_impact:.4f} & {r.win_rate*100:.1f} \\\\"
        )
    body_lines.extend([r"\bottomrule", r"\end{tabular}"])
    return "\n".join(body_lines)


def tca_table(runs: dict[str, dict]) -> str:
    rows = []
    for n in HEADLINE_RUNS:
        t = runs[n]["tca"]
        rows.append(
            {
                "config": n,
                "gross_pnl": float(t.get("gross_pnl", float("nan"))),
                "impact_cost": float(t.get("impact_cost", float("nan"))),
                "net_pnl": float(t.get("net_pnl", float("nan"))),
                "predicted_alpha": float(t.get("predicted_alpha", float("nan"))),
                "realised_alpha": float(t.get("realised_alpha", float("nan"))),
                "turnover": float(t.get("turnover", float("nan"))),
            }
        )
    df = pd.DataFrame(rows)
    body_lines = [
        r"\begin{tabular}{lrrrrrr}",
        r"\toprule",
        (r"Config & Gross P\&L & Impact cost & Net P\&L & "
         r"Predicted $\alpha$ & Realised $\alpha$ & Turnover \\"),
        r"\midrule",
    ]

    def fmt(v):
        return f"{v:,.0f}" if np.isfinite(v) else "--"

    for r in df.itertuples(index=False):
        cfg = r.config.replace("_", r"\_")
        body_lines.append(
            f"\\texttt{{{cfg}}} & {fmt(r.gross_pnl)} & {fmt(r.impact_cost)} & "
            f"{fmt(r.net_pnl)} & {fmt(r.predicted_alpha)} & "
            f"{fmt(r.realised_alpha)} & {fmt(r.turnover)} \\\\"
        )
    body_lines.extend([r"\bottomrule", r"\end{tabular}"])
    return "\n".join(body_lines)


# ---------------------------------------------------------------------------
# Sensitivity sweep -- run extra configs against the live pipeline.
# ---------------------------------------------------------------------------
def run_sensitivity(panel, alphas, lam_lookup_daily, H_star, TAU_BINS, RHO):
    """Stress sweeps: wrong impact model, signal delay, rho sweep, H sweep."""
    bins, daily_stats = panel.bins, panel.daily_stats
    rows: list[dict] = []

    # 1. Baseline (OW/linear, daily-reset) -- reuse already-saved metrics.
    base_metrics = pd.read_csv(SAVED / "ow_daily" / "metrics.csv", index_col=0)["value"]
    rows.append(
        {
            "scenario": "Baseline OW (headline)",
            "sharpe_sim": float(base_metrics["sharpe_sim"]),
            "net_pnl": float(base_metrics["total_pnl_sim"]),
            "impact_cost": float(base_metrics["total_impact_cost"]),
        }
    )

    # 2. Wrong model: trade against sqrt-impact while fit was linear.
    cfg = pi.BacktestConfig(
        name="sens_wrong_model", model_type="sqrt", strategy="afs",
        carry="daily", half_life_minutes=H_star, tau_bins=TAU_BINS,
        rho=RHO, save_root=SAVED,
    )
    out = pi.run_and_save(bins, daily_stats, alphas, lam_lookup_daily["linear"], cfg)
    rows.append(
        {
            "scenario": "Wrong model (fit OW, trade AFS)",
            "sharpe_sim": out.metrics["sharpe_sim"],
            "net_pnl": out.metrics["total_pnl_sim"],
            "impact_cost": out.metrics["total_impact_cost"],
        }
    )

    # 3. Signal delayed by 60s (6 bins).
    delayed = alphas.copy()
    delayed["alpha"] = (
        delayed.groupby(["stock", "date"])["alpha"].shift(6).fillna(0.0)
    )
    cfg = pi.BacktestConfig(
        name="sens_signal_delay_1min", model_type="linear", strategy="ow",
        carry="daily", half_life_minutes=H_star, tau_bins=TAU_BINS,
        rho=RHO, save_root=SAVED,
    )
    out = pi.run_and_save(bins, daily_stats, delayed, lam_lookup_daily["linear"], cfg)
    rows.append(
        {
            "scenario": "Signal delayed 1 min",
            "sharpe_sim": out.metrics["sharpe_sim"],
            "net_pnl": out.metrics["total_pnl_sim"],
            "impact_cost": out.metrics["total_impact_cost"],
        }
    )

    # 4. Rho sweep (signal strength).
    for rho_test in (0.01, 0.10):
        alphas_r = pi.create_synthetic_alpha(bins, rho=rho_test, h_bins=1, seed=42)
        cfg = pi.BacktestConfig(
            name=f"sens_rho_{int(rho_test*100):02d}", model_type="linear",
            strategy="ow", carry="daily", half_life_minutes=H_star,
            tau_bins=TAU_BINS, rho=rho_test, save_root=SAVED,
        )
        out = pi.run_and_save(bins, daily_stats, alphas_r, lam_lookup_daily["linear"], cfg)
        rows.append(
            {
                "scenario": rf"$\rho={rho_test:.2f}$",
                "sharpe_sim": out.metrics["sharpe_sim"],
                "net_pnl": out.metrics["total_pnl_sim"],
                "impact_cost": out.metrics["total_impact_cost"],
            }
        )

    # 5. Half-life sensitivity (60min baseline already covered; test 30 and 120).
    for H_test in (30.0, 120.0):
        cfg = pi.BacktestConfig(
            name=f"sens_H_{int(H_test)}", model_type="linear", strategy="ow",
            carry="daily", half_life_minutes=H_test, tau_bins=int(H_test * 6),
            rho=RHO, save_root=SAVED,
        )
        out = pi.run_and_save(bins, daily_stats, alphas, lam_lookup_daily["linear"], cfg)
        rows.append(
            {
                "scenario": f"Half-life = {int(H_test)} min",
                "sharpe_sim": out.metrics["sharpe_sim"],
                "net_pnl": out.metrics["total_pnl_sim"],
                "impact_cost": out.metrics["total_impact_cost"],
            }
        )

    df = pd.DataFrame(rows)
    body = [
        r"\begin{tabular}{lrrr}",
        r"\toprule",
        r"Scenario & Sharpe & Net P\&L (\$) & Impact cost (\$) \\",
        r"\midrule",
    ]
    for r in df.itertuples(index=False):
        body.append(
            f"{r.scenario} & {r.sharpe_sim:+.3f} & {r.net_pnl:,.0f} & {r.impact_cost:,.0f} \\\\"
        )
    body.extend([r"\bottomrule", r"\end{tabular}"])
    return "\n".join(body), df


# ---------------------------------------------------------------------------
# Slides.
# ---------------------------------------------------------------------------
def build_slides(runs: dict[str, dict], sensitivity_df: pd.DataFrame):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle):
        layout = prs.slide_layouts[6]  # blank
        s = prs.slides.add_slide(layout)
        tx = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(1.5))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(40)
            r.font.bold = True
        tx2 = s.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12), Inches(1.0))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        for r in p2.runs:
            r.font.size = Pt(22)
        return s

    def add_content_slide(title, bullets=None, image=None, image_w=8.0,
                          image_x=4.8, image_y=1.5):
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        # Title
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        for r in p.runs:
            r.font.size = Pt(28); r.font.bold = True
        # Bullets
        if bullets:
            box_w = 4.2 if image else 12.0
            tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(box_w), Inches(5.5))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = f"• {b}"
                for r in para.runs:
                    r.font.size = Pt(16)
        if image and Path(image).exists():
            s.shapes.add_picture(str(image), Inches(image_x), Inches(image_y),
                                  width=Inches(image_w))
        return s

    # ---- Title ----
    add_title_slide(
        "Quantitative Trading and Price Impact",
        "Final Project — MSc Mathematics & Finance, 2025-2026     |     Anran Severac",
    )

    # ---- Pipeline overview ----
    add_content_slide(
        "Pipeline at a glance",
        [
            "Modular package src/price_impact/ — 8 single-responsibility modules",
            "Notebook = thin orchestration; every run writes to saved/<name>/",
            "Baseline met for all 7 sections of the assignment",
            "Two enhancements emphasised:",
            "   (a) Non-parametric impact + cross-stock γ regularisation (§2.2)",
            "   (b) Multi-day carry backtest engine + split detection (§2.3)",
            "Daily-reset vs multi-day Ī produced side-by-side as two columns",
            "  on one panel; one BacktestConfig flag switches downstream",
            "λ(t) extended-OW plumbing already in simulator (placeholder strategy)",
        ],
    )

    # ---- Data ----
    add_content_slide(
        "§2.1 Data Preparation",
        [
            "Panel: 2019, 10-second bins, top-20 US equities by abs(orderFlow)",
            "Daily stats: trailing 20-day σ (10-s return std) and ADV (Σ|q|)",
            "OOS protocol: 10 rolling windows; train m, tune m+1, validate m+2",
            f"Panel size: {runs['ow_daily']['metrics'].get('n_stockdays', 'N/A')} stock-days",
            "build_panel(data_dir, year, top_n=20, lookback_days=20) → PanelData",
        ],
    )

    # ---- Impact-state recursion ----
    add_content_slide(
        "§2.2 Impact-state recursion — OW vs AFS",
        [
            "Ī_{t+1} = (1-β) Ī_t + q̃_t,  β = ln 2 / H_bins",
            "OW (linear):   q̃_t = σ · q_t / ADV",
            "AFS (sqrt):    q̃_t = σ · sign(q) · √(|q|/ADV)",
            "Same engine produces daily-reset and multi-day Ī as two columns",
            "compute_impact_states(...) → [stock, date, time, q_tilde,",
            "    I_bar_daily, I_bar_multi]",
        ],
    )

    # ---- Half-life ----
    add_content_slide(
        "§2.2 Half-life grid search → H* ≈ 52.5 min",
        [
            "9-point grid × 2 models",
            "OW peaks at H=60, AFS at H=45",
            "OOS R² flat in 45–60 min region (<0.001)",
            "Single universal H* = 52.5 min",
            "  adopted (midpoint, ≤0.0001 cost)",
        ],
        image=PROJECT_ROOT / "figures" / "halflife_grid_search.png",
        image_w=7.0, image_x=5.5,
    )

    # ---- Non-parametric ----
    add_content_slide(
        "§2.2 Non-parametric impact + γ shrinkage  [advanced]",
        [
            "B = 15 quantile bins on train month (pooled)",
            "Per-stock g_{i,b} shrunk to universal ḡ_b",
            "γ tuned on test month, OOS on m+2",
            "Apples-to-apples pooled R²:",
            "   OW:  0.161 → 0.221  (+5.9 pp — linearity binding)",
            "   AFS: 0.269 → 0.263  (square-root on efficient frontier)",
        ],
        image=PROJECT_ROOT / "figures" / "impact_curves.png",
        image_w=7.0, image_x=5.5,
    )

    # ---- Regularisation regime-correctness ----
    add_content_slide(
        "§2.2 Regularisation kicks in correctly under stress",
        [
            "Headline (~3000 obs/cell): γ tuning curves flat",
            "Shrink train window 21d → 5d → 2d, n drops ~10×",
            "γ*/n: 0.16 → 0.56 → 0.93",
            "OW gain (reg − raw): 0.0005 → 0.0070  (~14×)",
            "Estimator is regime-correct: flatness in headline is",
            "  a property of the data-rich regime, not a defect",
        ],
        image=PROJECT_ROOT / "figures" / "stress6b_tuning_progression.png",
        image_w=7.0, image_x=5.5,
    )

    # ---- Backtest engine ----
    add_content_slide(
        "§2.3 Backtest engine — Waelbroeck + carry modes  [advanced]",
        [
            "P_sim(t) = P_0 (1 + cumret(t) + g_sim(t) − g_ref(t))",
            "g(t) = λ·Ī(t)  or  λ(t)·Ī(t)  (extended-OW plumbing in place)",
            "Generic trade interface: trade_provider callable accepts any",
            "  signed schedule (optimal, TWAP, VWAP, manual)",
            "carry='daily' vs carry='multi'  (overnight decay 16h × β)",
            "Split detection helper: notional-stable price/volume jumps",
            "run_backtest(merged, model, trade_provider, daily_stats, carry)",
        ],
    )

    # ---- Synthetic alpha + strategy ----
    add_content_slide(
        "§2.4–§2.5 Synthetic alpha + optimal strategy",
        [
            "α_t = r^h_t + y · dW_t / P_t  with x=1 (unbiased forecast)",
            "y calibrated so Corr(α, r^h) = ρ;  ρ=0.05, h=1 bin",
            "OW closed form:  X*_t = α_t · ADV / (2λσ),  κ = β",
            "AFS: same target rule (per project.ipynb), sqrt enters in sim",
            "Extended OW with λ(t): placeholder (HJB derivation pending)",
            "End-of-day liquidation ramp (final 30 min)",
        ],
    )

    # ---- Cumulative P&L ----
    cum_pnl_path = SAVED / "ow_daily" / "cum_pnl.png"
    add_content_slide(
        "§2.6 Cumulative P&L — OW, daily-reset",
        [
            "Net P&L (Waelbroeck sim) vs gross P&L (mid)",
            "Gap = realised impact cost",
            f"Total net P&L: ${float(runs['ow_daily']['metrics']['total_pnl_sim']):,.0f}",
            f"Total impact cost: ${float(runs['ow_daily']['metrics']['total_impact_cost']):,.0f}",
            f"Annualised Sharpe (net): {float(runs['ow_daily']['metrics']['sharpe_sim']):+.3f}",
        ],
        image=cum_pnl_path, image_w=8.0, image_x=4.5, image_y=2.0,
    )

    # ---- Comparison table ----
    cmp_lines = ["Config           Days   Sharpe   Net P&L($)   ImpactCost($)   MaxDD($)    Win%"]
    for n in HEADLINE_RUNS:
        m = runs[n]["metrics"]
        cmp_lines.append(
            f"{n:14s}  {int(float(m['n_days'])):>4}  "
            f"{float(m['sharpe_sim']):+.3f}  "
            f"{float(m['total_pnl_sim']):>10,.0f}  "
            f"{float(m['total_impact_cost']):>13,.0f}  "
            f"{float(m['max_drawdown']):>9,.0f}  "
            f"{float(m['win_rate'])*100:>4.1f}"
        )
    add_content_slide(
        "§2.6 Headline comparison (4 configs)", cmp_lines,
    )

    # ---- TCA ----
    tca_lines = ["Config           Gross($)    Impact($)    Net($)    PredAlpha($)  ReadAlpha($)"]
    for n in HEADLINE_RUNS:
        t = runs[n]["tca"]
        def fmt(v):
            return f"{float(v):,.0f}" if pd.notna(v) else "--"
        tca_lines.append(
            f"{n:14s}  {fmt(t.get('gross_pnl')):>10s}  "
            f"{fmt(t.get('impact_cost')):>10s}  {fmt(t.get('net_pnl')):>9s}  "
            f"{fmt(t.get('predicted_alpha')):>11s}  {fmt(t.get('realised_alpha')):>11s}"
        )
    add_content_slide(
        "§2.6 TCA decomposition (totals)", tca_lines,
    )

    # ---- Sensitivity ----
    sens_lines = ["Scenario                              Sharpe    Net P&L($)    Impact cost($)"]
    for _, r in sensitivity_df.iterrows():
        sens_lines.append(
            f"{r['scenario'][:36]:36s}  {r['sharpe_sim']:+.3f}  "
            f"{r['net_pnl']:>10,.0f}  {r['impact_cost']:>13,.0f}"
        )
    add_content_slide(
        "§2.7 Sensitivity & stress tests", sens_lines,
    )

    # ---- Conclusion ----
    add_content_slide(
        "What we learned",
        [
            "Square-root impact (AFS) is on the efficient frontier for this panel;",
            "  non-parametric fit ties AFS, beats OW by +5.9 pp pooled OOS R²",
            "Cross-stock γ shrinkage is regime-correct: immaterial when data-rich,",
            "  γ*/n rises 6× when training window shrinks 10×",
            "Multi-day carry changes net P&L and impact cost vs daily reset",
            "  — visible in the headline comparison table",
            "Engine is modular: 4 configs ran via one for-loop over BacktestConfig;",
            "  artifacts auto-saved to saved/<name>/",
        ],
    )

    add_content_slide(
        "Thank you",
        ["Questions?"],
    )

    out_path = ROOT / "slides.pptx"
    prs.save(str(out_path))
    print(f"slides written: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Entry point.
# ---------------------------------------------------------------------------
def main():
    print("Loading run artifacts...")
    runs = {n: load_run(n) for n in HEADLINE_RUNS}

    print("Building headline metrics table...")
    (TABLES / "headline_metrics.tex").write_text(headline_table(runs))

    print("Building TCA summary table...")
    (TABLES / "tca_summary.tex").write_text(tca_table(runs))

    print("Building panel + lambda lookup for sensitivity sweep...")
    panel = pi.build_panel(PROJECT_ROOT / "data", year=2019, top_n=20)
    alphas = pi.create_synthetic_alpha(panel.bins, rho=0.05, h_bins=1, seed=42)
    # Reconstruct per-stock λ via a single OLS fit on the full year
    # (cheaper than re-running the rolling baseline here).
    impact_lin = pi.compute_impact_states(panel.bins, panel.daily_stats,
                                           half_life_minutes=52.5, model_type="linear")
    feats_lin = pi.build_regression_features(impact_lin, panel.bins,
                                              tau_bins=180, carry="daily")
    stats_lin = pi.daily_sufficient_stats(feats_lin)
    base = pi.rolling_baseline(stats_lin, n_windows=10, offset=2)
    lam_lookup_daily = {"linear": pi.per_stock_lambda(base).to_dict()}

    print("Running sensitivity sweep...")
    sens_body, sens_df = run_sensitivity(panel, alphas, lam_lookup_daily,
                                          H_star=52.5, TAU_BINS=180, RHO=0.05)
    (TABLES / "sensitivity.tex").write_text(sens_body)

    print("Building slides.pptx...")
    build_slides(runs, sens_df)

    print("Done.  Artifacts in", TABLES, "and", ROOT / "slides.pptx")


if __name__ == "__main__":
    main()
